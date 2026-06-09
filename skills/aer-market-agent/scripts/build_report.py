#!/usr/bin/env python3
"""Build the Default Market Offer (DMO) PowerPoint report from data/dmo_dataset.json.

    pip install python-pptx matplotlib openpyxl
    python scripts/build_report.py --out "DMO Report.pptx"

All content lives in data/dmo_dataset.json - edit that to change numbers, text or
branding. Add your own slide in add_custom_analysis() (clearly marked, near the end).

Built by Renovara - renovara.co
"""
import argparse, json, os, tempfile
import matplotlib; matplotlib.use("Agg")
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

DISPLAY, BODY = "Calibri Light", "Calibri"
HERE = os.path.dirname(os.path.abspath(__file__))
def C(h): return RGBColor.from_string(h)

def text(s,t,l,tp,w,h,size=14,color="0A1F44",bold=False,font=BODY,align=PP_ALIGN.LEFT,anchor=None,italic=False,spacing=None,link=None):
    tb=s.shapes.add_textbox(Inches(l),Inches(tp),Inches(w),Inches(h)); tf=tb.text_frame; tf.word_wrap=True
    for m in ("margin_left","margin_right","margin_top","margin_bottom"): setattr(tf,m,0)
    if anchor: tf.vertical_anchor=anchor
    p=tf.paragraphs[0]; p.alignment=align
    if spacing: p.line_spacing=spacing
    r=p.add_run(); r.text=t; f=r.font; f.size=Pt(size); f.bold=bold; f.italic=italic; f.name=font; f.color.rgb=C(color)
    if link:
        try: r.hyperlink.address=link
        except Exception: pass
    return tb
def rect(s,l,tp,w,h,color,line=None):
    sp=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(l),Inches(tp),Inches(w),Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb=C(color)
    if line: sp.line.color.rgb=C(line); sp.line.width=Pt(1)
    else: sp.line.fill.background()
    sp.shadow.inherit=False; return sp
def oval(s,l,tp,w,h,color):
    sp=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(l),Inches(tp),Inches(w),Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb=C(color); sp.line.fill.background(); sp.shadow.inherit=False; return sp

def _ax(d,w,h):
    fig,ax=plt.subplots(figsize=(w,h),dpi=200); fig.patch.set_facecolor("#"+d["branding"]["paper"]); ax.set_facecolor("#"+d["branding"]["paper"]); return fig,ax
def _fin(fig,tmp,name):
    p=os.path.join(tmp,name); fig.tight_layout(); fig.savefig(p,facecolor=fig.get_facecolor()); plt.close(fig); return p

def chart_longview(d,tmp):
    v=d["longview"]["revenue"]; xs=list(range(2006,2006+len(v))); fig,ax=_ax(d,7.2,4.3)
    ax.plot(xs,v,color="#"+d["branding"]["text"],linewidth=2.6,marker="o",markersize=4)
    ax.set_ylim(0,16); ax.spines[["top","right"]].set_visible(False); ax.spines[["left","bottom"]].set_color("#cdd5e0")
    ax.tick_params(colors="#5B6378",labelsize=8); ax.yaxis.set_major_formatter(lambda y,_:f"${y:,.0f}"); ax.grid(axis="y",color="#D9DEE7",linewidth=0.6)
    return _fin(fig,tmp,"lv.png")
def chart_usage(d,tmp):
    v=d["usage_per_customer"]["values"]; xs=list(range(2006,2006+len(v))); fig,ax=_ax(d,6.8,4.0)
    ax.plot(xs,v,color="#"+d["branding"]["text"],linewidth=2.6,marker="o",markersize=4)
    ax.set_ylim(10,18); ax.spines[["top","right"]].set_visible(False); ax.spines[["left","bottom"]].set_color("#cdd5e0")
    ax.tick_params(colors="#5B6378",labelsize=8); ax.grid(axis="y",color="#D9DEE7",linewidth=0.6)
    return _fin(fig,tmp,"usage.png")
def chart_traj(d,tmp):
    fys=[x[1][2:] for x in d["determinations"]]; fig,ax=_ax(d,7.6,4.3); mk=["o","s","^","D","P"]
    for i,z in enumerate(d["zones"]):
        ax.plot(fys,d["residential_flat_aud"][z["id"]],color="#"+z["color"],linewidth=2.4,marker=mk[i%5],markersize=5,label=z["label"])
    ax.set_ylim(1200,2900); ax.spines[["top","right"]].set_visible(False); ax.spines[["left","bottom"]].set_color("#cdd5e0")
    ax.tick_params(colors="#5B6378",labelsize=9); ax.yaxis.set_major_formatter(lambda v,_:f"${v:,.0f}")
    ax.grid(axis="y",color="#D9DEE7",linewidth=0.6); ax.legend(fontsize=8,frameon=False,ncol=2,loc="upper left")
    return _fin(fig,tmp,"traj.png")
def chart_realnom(d,tmp):
    fys=[x[1][2:] for x in d["determinations"]]; fig,ax=_ax(d,7.0,4.3)
    ax.plot(fys,d["residential_flat_aud"]["ausgrid"],color="#0A1F44",linewidth=2.6,marker="o",markersize=5,label="Nominal (as billed)")
    ax.plot(fys,d["residential_real_ausgrid"],color="#00C2A8",linewidth=2.6,marker="s",markersize=5,label="Real (2019-20 dollars)")
    ax.set_ylim(1200,2100); ax.spines[["top","right"]].set_visible(False); ax.spines[["left","bottom"]].set_color("#cdd5e0")
    ax.tick_params(colors="#5B6378",labelsize=9); ax.yaxis.set_major_formatter(lambda v,_:f"${v:,.0f}")
    ax.grid(axis="y",color="#D9DEE7",linewidth=0.6); ax.legend(fontsize=9,frameon=False,ncol=2,loc="upper left")
    return _fin(fig,tmp,"realnom.png")
def chart_doughnut(d,tmp):
    labels=list(d["cost_stack_proportions"].keys()); vals=[d["cost_stack_proportions"][k] for k in labels]; cols=["#"+d["cost_stack_colors"][k] for k in labels]
    fig,ax=_ax(d,4.2,4.2)
    w,_,at=ax.pie(vals,colors=cols,startangle=90,counterclock=False,wedgeprops=dict(width=0.42,edgecolor="#"+d["branding"]["paper"],linewidth=2),autopct=lambda p:f"{p:.0f}%",pctdistance=0.79,textprops=dict(fontsize=11,weight="bold"))
    for i,t in enumerate(at): t.set_color("#0A1F44" if labels[i]=="Environmental" else "white")
    ax.text(0,0.07,f"{vals[0]}%",ha="center",va="center",fontsize=22,weight="bold",color=cols[0]); ax.text(0,-0.16,"is network",ha="center",va="center",fontsize=10,color="#5B6378")
    return _fin(fig,tmp,"dough.png")
def chart_components(d,tmp):
    ts=d["cost_stack_timeseries"]; fys=ts["fy"]; order=["Network","Wholesale","Retail + margin","Environmental"]; fig,ax=_ax(d,7.4,4.5); bottom=[0]*len(fys)
    for k in order:
        ax.bar(fys,ts[k],bottom=bottom,color="#"+d["cost_stack_colors"][k],label=k,width=0.66); bottom=[a+b for a,b in zip(bottom,ts[k])]
    ax.set_ylim(0,2050); ax.spines[["top","right"]].set_visible(False); ax.spines[["left","bottom"]].set_color("#cdd5e0")
    ax.tick_params(colors="#5B6378",labelsize=9); ax.yaxis.set_major_formatter(lambda v,_:f"${v:,.0f}")
    ax.grid(axis="y",color="#D9DEE7",linewidth=0.6); ax.legend(fontsize=8,frameon=False,ncol=4,loc="upper center",bbox_to_anchor=(0.5,-0.08))
    return _fin(fig,tmp,"comp.png")
def chart_demand(d,tmp):
    L=d["demand_labels"]; fig,ax=_ax(d,6.8,4.0)
    ax.bar(L,d["demand_avg"],color="#"+d["cost_stack_colors"]["Wholesale"],width=0.6,label="Average demand")
    ax.plot(L,d["demand_peak"],color="#"+d["branding"]["text"],linewidth=3,marker="o",markersize=5,label="Peak demand")
    ax.set_ylim(0,40); ax.spines[["top","right"]].set_visible(False); ax.spines[["left","bottom"]].set_color("#cdd5e0")
    ax.tick_params(colors="#5B6378",labelsize=9); ax.grid(axis="y",color="#D9DEE7",linewidth=0.6)
    ax.legend(fontsize=9,frameon=False,ncol=2,loc="upper center",bbox_to_anchor=(0.5,-0.08))
    return _fin(fig,tmp,"dem.png")
def chart_regions(d,tmp):
    rows=d["regional_change_pct"]; labels=[r[0] for r in rows]; vals=[r[1] for r in rows]; teal="#"+d["branding"]["teal"]; gold="#"+d["branding"]["gold"]
    cols=[teal if v<0 else gold for v in vals]; fig,ax=_ax(d,6.4,4.3); y=list(range(len(labels)))
    ax.barh(y,[abs(v) for v in vals],color=cols,height=0.6); ax.set_yticks(y); ax.set_yticklabels(labels,fontsize=10,color="#0A1F44"); ax.invert_yaxis()
    for i,v in enumerate(vals): ax.text(abs(v)+0.12,i,("+" if v>0 else "-")+f"{abs(v):.1f}%",va="center",fontsize=11,weight="bold",color=("#008C7A" if v<0 else "#9A7D00"))
    ax.set_xlim(0,max(abs(v) for v in vals)+1.3); ax.spines[["top","right","bottom"]].set_visible(False); ax.spines["left"].set_color("#cdd5e0"); ax.tick_params(axis="x",colors="#5B6378",labelsize=9)
    return _fin(fig,tmp,"reg.png")

def blank(prs): return prs.slides.add_slide(prs.slide_layouts[6])
def bg(s,prs,c): rect(s,0,0,prs._w,prs._h,c)
def acc(s,prs,c,w=0.16): rect(s,0,0,w,prs._h,c)
def eyebrow(s,t,c,l=0.7,tp=0.5): text(s,t.upper(),l,tp,11,0.3,11,c,True)
def srcfoot(s,prs,src,dark=False):
    if src: text(s,src,prs._w-8.2,prs._h-0.4,8.0,0.3,8.5,(prs._b["inkmute"] if dark else prs._b["muted"]),italic=True,align=PP_ALIGN.RIGHT)
def tbl(s,x,y,cols,rows,headfill,fs=10.5,rh=0.34):
    # simple manual table: cols=[(label,width)], rows=list of list of (text,color,bold)
    cx=x
    for ci,(lab,w) in enumerate(cols):
        rect(s,cx,y,w,rh,headfill); text(s,lab,cx+0.08,y,w-0.12,rh,fs,"FFFFFF",True,anchor=MSO_ANCHOR.MIDDLE); cx+=w
    for ri,row in enumerate(rows):
        ry=y+rh*(ri+1); cx=x
        for ci,(lab,w) in enumerate(cols):
            rect(s,cx,ry,w,rh,"FFFFFF",line="D9DEE7"); cell=row[ci]
            text(s,cell[0],cx+0.08,ry,w-0.12,rh,fs,cell[1] if len(cell)>1 else "1B2A3A",cell[2] if len(cell)>2 else False,anchor=MSO_ANCHOR.MIDDLE); cx+=w

def build(d,out):
    prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5); prs._w,prs._h,prs._b=13.333,7.5,d["branding"]
    b=d["branding"]; NAVY,TEAL,TEALD,GOLD,SLATE,INK,INK2,PAPER,LINE,INKMUTE=b["text"],b["teal"],b["teald"],b["gold"],b["muted"],b["ink"],b["ink2"],b["paper"],b["line"],b["inkmute"]
    cc=d["cost_stack_colors"]; tmp=tempfile.mkdtemp()
    g_traj,g_dough,g_comp,g_real,g_dem,g_reg,g_lv,g_usage=chart_traj(d,tmp),chart_doughnut(d,tmp),chart_components(d,tmp),chart_realnom(d,tmp),chart_demand(d,tmp),chart_regions(d,tmp),chart_longview(d,tmp),chart_usage(d,tmp)

    # 1 TITLE
    s=blank(prs); bg(s,prs,INK); acc(s,prs,TEAL,0.18)
    eyebrow(s,b["org"]+"  -  Australian energy market analysis",TEAL,0.95,0.95)
    text(s,"Where your power dollar goes",0.95,1.7,11.5,2.0,56,"FFFFFF",False,DISPLAY,spacing=1.0)
    text(s,"Seven years of Australia's Default Market Offer, decomposed",0.97,3.9,11,0.6,22,"D9DEE7")
    text(s,"NSW   -   South-East QLD   -   South Australia        DMO 1-8   -   2019-2027",0.97,4.7,11.5,0.4,15,INKMUTE)
    text(s,"Source: "+d["meta"]["source"],0.95,prs._h-0.5,11,0.3,10,INKMUTE,italic=True)

    # 2 BASICS
    s=blank(prs); bg(s,prs,PAPER); acc(s,prs,TEAL); eyebrow(s,"The basics",TEALD)
    text(s,"A safety net - and the market's yardstick",0.7,0.85,12,0.8,32,NAVY,True)
    pal=[TEAL,NAVY,GOLD,SLATE]
    for i,row in enumerate(d["basics"]):
        yy=2.15+i*1.18; oval(s,0.95,yy,0.5,0.5,pal[i%4]); text(s,row[0],1.7,yy-0.05,10.5,0.4,17,NAVY,True); text(s,row[1],1.7,yy+0.42,10.6,0.6,13.5,SLATE)

    # 3 HOW CALCULATED
    s=blank(prs); bg(s,prs,PAPER); acc(s,prs,TEAL); eyebrow(s,"How the numbers are derived",TEALD)
    text(s,"How the AER builds a DMO bill",0.7,0.85,12,0.8,32,NAVY,True)
    text(s,"The AER builds the price bottom-up for a model customer, then adds GST:",0.7,1.95,7,0.4,13.5,SLATE,italic=True)
    ccol=[cc["Network"],cc["Wholesale"],cc["Environmental"],SLATE,NAVY]
    for i,r in enumerate(d["calc_components"]):
        yy=2.55+i*0.82; rect(s,0.7,yy+0.04,0.28,0.28,ccol[i])
        text(s,r[0]+"  -  "+r[1],1.1,yy-0.04,6.7,0.6,12.5,NAVY,False,anchor=MSO_ANCHOR.MIDDLE)
    text(s,"(  sum of components  )  x  GST  =  annual reference price",0.7,6.7,7.5,0.4,13,TEALD,True)
    we=d["worked_example"]; rect(s,8.5,2.0,4.3,4.6,"FFFFFF",line=LINE); rect(s,8.5,2.0,4.3,0.7,NAVY)
    text(s,we["title"],8.75,2.12,3.9,0.45,14,"FFFFFF",True,anchor=MSO_ANCHOR.MIDDLE)
    text(s,we["usage"],8.75,2.8,3.9,0.3,10.5,SLATE,italic=True)
    wc=[cc["Network"],cc["Wholesale"],SLATE,cc["Environmental"]]
    for i,r in enumerate(we["rows"]):
        yy=3.3+i*0.55; rect(s,8.75,yy+0.06,0.2,0.2,wc[i]); text(s,r[0],9.05,yy,2.6,0.4,13,NAVY,anchor=MSO_ANCHOR.MIDDLE); text(s,"$"+str(r[1]),11.4,yy,1.2,0.4,13,NAVY,True,align=PP_ALIGN.RIGHT,anchor=MSO_ANCHOR.MIDDLE)
    rect(s,8.75,5.6,3.85,0.02,LINE); text(s,"Reference price",9.05,5.7,2.4,0.45,14,NAVY,True,anchor=MSO_ANCHOR.MIDDLE); text(s,"$"+f"{we['total']:,}",11.0,5.68,1.6,0.5,19,TEALD,True,align=PP_ALIGN.RIGHT,anchor=MSO_ANCHOR.MIDDLE)
    text(s,we["note"],8.75,6.25,3.85,0.35,8.5,SLATE,italic=True)
    srcfoot(s,prs,"AER DMO methodology - ACIL Allen wholesale modelling")

    # 4 TRAJECTORY
    s=blank(prs); bg(s,prs,PAPER); acc(s,prs,TEAL); eyebrow(s,"The seven-year arc",TEALD)
    text(s,"The bill went on a rollercoaster",0.7,0.85,12,0.8,32,NAVY,True)
    s.shapes.add_picture(g_traj,Inches(0.6),Inches(2.0),height=Inches(4.8))
    rect(s,9.7,2.2,3.0,4.0,"FFFFFF",line=LINE); rect(s,9.7,2.2,0.09,4.0,GOLD)
    text(s,"+36% to +50%",9.95,2.45,2.6,0.6,24,NAVY,True); text(s,"trough (2021) to peak",9.95,3.1,2.6,0.4,12,SLATE)
    text(s,"Then 2026 brings the first real relief - except in South Australia.",9.95,3.7,2.65,1.2,14,NAVY)
    text(s,"Outside regional NSW, SA is now the most expensive zone.",9.95,5.05,2.65,1.0,11,TEALD,italic=True)
    srcfoot(s,prs,"AER final determinations, DMO 1-8 - residential flat rate")

    # 5 ANATOMY
    s=blank(prs); bg(s,prs,PAPER); acc(s,prs,TEAL); eyebrow(s,"Where your dollar goes",TEALD)
    text(s,"Four costs stack up to your bill",0.7,0.85,12,0.8,32,NAVY,True)
    s.shapes.add_picture(g_dough,Inches(0.7),Inches(2.1),height=Inches(4.6))
    for i,k in enumerate(d["cost_stack_proportions"].keys()):
        cy=2.2+i*1.15; rect(s,6.6,cy+0.05,0.28,0.28,cc[k]); text(s,f"{k}   {d['cost_stack_proportions'][k]}%",7.05,cy-0.05,5.6,0.4,16,NAVY,True); text(s,d["cost_stack_notes"][k],7.05,cy+0.38,5.6,0.7,12.5,SLATE)
    srcfoot(s,prs,"Indicative split - AER DMO 8 cost-stack proportions; varies by zone")

    # 6 COMPONENTS OVER TIME
    s=blank(prs); bg(s,prs,PAPER); acc(s,prs,TEAL); eyebrow(s,"Where your dollar goes over time",TEALD)
    text(s,"Every component, tracked year by year",0.7,0.85,12,0.8,32,NAVY,True)
    s.shapes.add_picture(g_comp,Inches(0.5),Inches(1.9),height=Inches(4.6))
    text(s,"Year-on-year, 2025-26 to 2026-27",8.7,1.95,4.2,0.35,12,SLATE,True)
    yy_=d["cost_stack_timeseries"]["yoy_latest"]
    for i,k in enumerate(["Wholesale","Network","Retail + margin","Environmental"]):
        v=yy_[k]; yy=2.5+i*0.92; dn=v<0; rect(s,8.7,yy,4.1,0.76,"FFFFFF",line=LINE); rect(s,8.7,yy,0.09,0.76,cc[k])
        text(s,k,8.95,yy+0.12,2.5,0.5,14,NAVY,True,anchor=MSO_ANCHOR.MIDDLE)
        text(s,("v " if dn else "^ ")+("-" if dn else "+")+f"{abs(v):.1f}%",11.2,yy+0.12,1.45,0.5,15,(TEALD if dn else "9A7D00"),True,align=PP_ALIGN.RIGHT,anchor=MSO_ANCHOR.MIDDLE)
    text(s,"Nominal dollars; adjusted for inflation the rise is far smaller (next slide). Bars reconcile to the published bill.",0.5,6.7,8.0,0.3,9.5,SLATE,italic=True)
    srcfoot(s,prs,"Ausgrid residential - wholesale & environmental from AER cost data; network/retail split indicative")

    # 7 REAL VS NOMINAL
    s=blank(prs); bg(s,prs,PAPER); acc(s,prs,TEAL); eyebrow(s,"Adjusted for inflation",TEALD)
    text(s,"Adjust for inflation, and the rise shrinks",0.7,0.85,12,0.8,32,NAVY,True)
    s.shapes.add_picture(g_real,Inches(0.55),Inches(2.0),height=Inches(4.5))
    text(s,"Ausgrid residential - nominal vs real",0.6,6.6,7,0.3,10,SLATE,italic=True)
    text(s,"2019-20 to 2026-27 change",8.5,1.95,4.3,0.35,13,NAVY,True)
    tbl(s,8.5,2.4,[("Zone",2.2),("Nominal",1.05),("Real",1.05)],[[(r[0],),(r[1],),(r[2],TEALD,True)] for r in d["nominal_vs_real_change"]],NAVY,fs=10.5,rh=0.36)
    rect(s,8.5,5.0,4.3,1.4,"E6F7F4"); rect(s,8.5,5.0,0.09,1.4,TEAL)
    text(s,"When we account for economy-wide inflation, the typical bill is roughly back to 2019 - SA and QLD are lower. The 2022 crisis spike was real, but has largely unwound.",8.75,5.06,3.9,1.32,11,NAVY,anchor=MSO_ANCHOR.MIDDLE,spacing=1.04)
    srcfoot(s,prs,d["cpi_note"])

    # 8 NETWORK LONG VIEW
    s=blank(prs); bg(s,prs,PAPER); acc(s,prs,TEAL); eyebrow(s,"Network - the long view",TEALD)
    text(s,"In real terms, the network peaked a decade ago",0.7,0.85,12,0.8,30,NAVY,True)
    s.shapes.add_picture(g_lv,Inches(0.55),Inches(1.95),height=Inches(4.4))
    scol=[NAVY,cc["Wholesale"],SLATE,GOLD]
    for i,c in enumerate(d["longview"]["stats"]):
        yy=2.0+i*1.12; rect(s,8.3,yy,4.5,0.98,"FFFFFF",line=LINE); rect(s,8.3,yy,0.09,0.98,scol[i])
        text(s,c[0],8.55,yy+0.08,2.3,0.3,12,SLATE,True)
        text(s,c[1],10.6,yy+0.05,2.1,0.4,18,(scol[i] if scol[i]!=GOLD else "9A7D00"),True,DISPLAY,align=PP_ALIGN.RIGHT)
        text(s,c[2],8.55,yy+0.47,4.1,0.4,10,SLATE)
    text(s,d["longview"]["insight"],0.7,6.5,12,0.45,11,SLATE,italic=True,spacing=1.0)
    srcfoot(s,prs,"AER State of the Energy Market 2025, ch.3 - real (June-2024 $) - NEM distribution aggregate")

    # 9 BUILDING BLOCKS
    s=blank(prs); bg(s,prs,PAPER); acc(s,prs,TEAL); eyebrow(s,"Network - decomposed",TEALD)
    text(s,"Network costs, decomposed",0.7,0.85,12,0.8,32,NAVY,True)
    text(s,"How allowed network revenue is built (approx. shares):",0.7,1.95,12,0.3,14,NAVY,True)
    bbx=0.7; bbw=12.0; bbcol=[cc["Network"],cc["Wholesale"],SLATE,GOLD]
    for i,seg in enumerate(d["network_building_blocks"]):
        w=bbw*seg[1]/100.0; rect(s,bbx,2.45,w,0.85,bbcol[i])
        if seg[1]>=12: text(s,f"{seg[0]}  {seg[1]}%",bbx,2.45,w,0.85,13,("0A1F44" if bbcol[i]==GOLD else "FFFFFF"),True,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
        bbx+=w
    text(s,"Depreciation 6%   -   Adjustments + tax 7%",8.7,3.4,4.0,0.3,10,SLATE,italic=True,align=PP_ALIGN.RIGHT)
    text(s,"Return on capital = the regulated asset base x the AER-allowed rate of return. It is the biggest block by far - which is why a rate-of-return reset or asset-base growth moves network costs more than anything else.",0.7,3.9,12,0.7,14,NAVY,spacing=1.12)
    rect(s,0.7,5.2,12.1,1.0,"E6F7F4"); rect(s,0.7,5.2,0.09,1.0,TEAL)
    text(s,d["network_attribution"],0.95,5.3,11.7,0.8,13.5,NAVY,anchor=MSO_ANCHOR.MIDDLE)
    srcfoot(s,prs,"Building-block shares: Ausgrid 2024-29 draft decision (indicative); final differs slightly")

    # 10 EVERY NETWORK
    s=blank(prs); bg(s,prs,PAPER); acc(s,prs,TEAL); eyebrow(s,"Network - every reset",TEALD)
    text(s,"Every DMO-region network just reset",0.7,0.85,12,0.8,32,NAVY,True)
    text(s,"Five-year allowed revenue the AER set for each network serving the DMO regions:",0.7,1.95,12,0.3,13.5,SLATE,italic=True)
    tbl(s,0.7,2.5,[("Distribution",3.0),("Period",1.2),("Revenue",1.4),("vs prior",1.4)],[[(r[0],),(r[1],),(r[2],),(r[3],"9A7D00" if r[3].startswith("+") else SLATE,r[3].startswith("+"))] for r in d["networks_distribution"]],NAVY,fs=11,rh=0.4)
    tbl(s,8.0,2.5,[("Transmission charges",2.7),("Period",1.0),("Revenue",1.1)],[[(r[0],),(r[1],),(r[2],)] for r in d["networks_transmission"]],TEALD,fs=11,rh=0.4)
    text(s,"Project EnergyConnect: $1.8bn -> $3.6bn",8.0,5.0,4.8,0.3,10,SLATE,italic=True)
    text(s,"The common thread: every reset in 2022-25 landed in a high-rate, high-inflation window - the AER attributes ~46-50% of the distribution increases to higher inflation and interest rates.",0.7,5.6,12,0.6,12.5,NAVY,italic=True,spacing=1.1)
    text(s,"* Ausgrid draft-decision figure; 'n/p' = % vs prior not published. Transmission is ~4-11% of a typical bill; Project EnergyConnect is a near-term cost but a net saving.",0.7,6.7,12,0.25,8.5,SLATE,italic=True)

    # 11 NETWORK VS DEMAND
    s=blank(prs); bg(s,prs,PAPER); acc(s,prs,TEAL); eyebrow(s,"Network - cost vs demand",TEALD)
    text(s,"Rising cost, flat grid",0.7,0.85,12,0.8,32,NAVY,True)
    s.shapes.add_picture(g_dem,Inches(0.55),Inches(2.0),height=Inches(4.0))
    text(s,"NEM, financial year - demand in GW. Peak = annual maximum (summer or winter).",0.6,6.1,7.5,0.3,10,SLATE,italic=True)
    dc=[NAVY,cc["Wholesale"],GOLD]
    for i,c in enumerate(d["demand_stats"]):
        yy=2.0+i*1.4; rect(s,8.3,yy,4.5,1.25,"FFFFFF",line=LINE); rect(s,8.3,yy,0.09,1.25,dc[i])
        text(s,c[0],8.55,yy+0.1,4.1,0.3,12,SLATE,True); text(s,c[1],8.55,yy+0.4,4.1,0.45,22,(dc[i] if dc[i]!=GOLD else "9A7D00"),True,DISPLAY); text(s,c[2],8.55,yy+0.92,4.05,0.3,9.5,SLATE)
    text(s,"Demand has been flat on both measures since 2019, yet allowed network revenue rose 33-47% - about half from higher financing costs (the rate-of-return reset and inflation) and about half from real investment in replacing and modernising an ageing grid, not load growth. (See next slide.)",0.55,6.45,12.2,0.5,9.5,SLATE,italic=True,spacing=1.0)
    srcfoot(s,prs,"AER demand & energy charts (NEM) - network revenue: AER determinations")

    # 11b WHY NETWORK BILL ROSE
    s=blank(prs); bg(s,prs,PAPER); acc(s,prs,TEAL); eyebrow(s,"Network - what's really driving it",TEALD)
    text(s,"Why the network bill rose — not demand",0.7,0.85,12,0.8,30,NAVY,True)
    text(s,d["why_caption"],0.7,1.85,12,0.3,13,SLATE,italic=True)
    fw=12.0; sx=0.7
    rect(s,sx,2.3,fw*0.46,0.85,NAVY); rect(s,sx+fw*0.46,2.3,fw*0.54,0.85,TEAL)
    text(s,"Financing + inflation  ~46%",sx,2.3,fw*0.46,0.85,14,"FFFFFF",True,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    text(s,"Real investment + opex  ~54%",sx+fw*0.46,2.3,fw*0.54,0.85,14,"0A1F44",True,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    text(s,"rate-of-return reset + the asset base re-priced by inflation",sx,3.2,fw*0.46,0.3,10,SLATE,italic=True,align=PP_ALIGN.CENTER)
    text(s,"replacing & modernising an ageing grid — not demand growth",sx+fw*0.46,3.2,fw*0.54,0.3,10,SLATE,italic=True,align=PP_ALIGN.CENTER)
    text(s,"What the investment is for",0.7,3.95,6,0.3,14,NAVY,True)
    icol=[cc["Network"],TEAL,GOLD]
    for i,r in enumerate(d["investment_for"]):
        yy=4.4+i*0.66; rect(s,0.7,yy+0.05,0.22,0.22,icol[i])
        text(s,r[0],1.05,yy-0.02,6.2,0.35,13,NAVY,True); text(s,"— "+r[1],1.05,yy+0.32,6.4,0.3,10.5,SLATE)
    text(s,"And the AER cut demand-driven augmentation — on lower demand forecasts.",0.7,6.6,7,0.3,11,TEALD,italic=True)
    rect(s,8.2,3.95,4.6,2.45,"FFFFFF",line=LINE); rect(s,8.2,3.95,0.09,2.45,GOLD)
    text(s,"Renewal, not a new build-out",8.45,4.1,4.2,0.3,13,TEALD,True)
    text(s,d["km_note"],8.45,4.5,4.25,1.85,11.5,NAVY,spacing=1.08)
    srcfoot(s,prs,"AER SA Power Networks 2025-30 final decision - network length from company disclosures")

    # 12 REGIONS
    s=blank(prs); bg(s,prs,PAPER); acc(s,prs,TEAL); eyebrow(s,"Regions",TEALD)
    text(s,"Same year, opposite directions",0.7,0.85,12,0.8,32,NAVY,True)
    s.shapes.add_picture(g_reg,Inches(1.2),Inches(2.0),height=Inches(4.6))
    text(s,"SA's network costs outweigh its wholesale savings - the lone riser.",1.4,6.55,10,0.4,13,NAVY,italic=True)
    srcfoot(s,prs,"Residential flat rate - DMO 8 (2026-27) vs DMO 7")

    # 12b TARIFF STRUCTURE SHIFT
    s=blank(prs); bg(s,prs,PAPER); acc(s,prs,TEAL); eyebrow(s,"Network - the structural shift",TEALD)
    text(s,"Falling usage is reshaping your bill",0.7,0.85,12,0.8,32,NAVY,True)
    s.shapes.add_picture(g_usage,Inches(0.55),Inches(2.0),height=Inches(4.0))
    text(s,"Grid energy used per distribution customer (MWh/yr)",0.6,6.1,7,0.3,10,SLATE,italic=True)
    ts=d["tariff_shift"]
    rect(s,8.3,2.0,4.5,1.2,"FFFFFF",line=LINE); rect(s,8.3,2.0,0.09,1.2,NAVY)
    text(s,ts["usage_value"],8.55,2.12,4.1,0.5,24,NAVY,True,DISPLAY)
    text(s,ts["usage_detail"],8.55,2.7,4.15,0.4,10,SLATE)
    rect(s,8.3,3.4,4.5,1.2,"FFFFFF",line=LINE); rect(s,8.3,3.4,0.09,1.2,TEAL)
    text(s,ts["tariff_value"],8.55,3.52,4.1,0.5,24,TEALD,True,DISPLAY)
    text(s,ts["tariff_detail"],8.55,4.1,4.15,0.4,10,SLATE)
    text(s,ts["text"],8.3,4.8,4.5,1.6,10.5,NAVY,spacing=1.05)
    srcfoot(s,prs,"AER State of the Energy Market 2025, ch.3 (usage Fig 3.31; tariffs Fig 3.9)")

    # 13 REFORMS (dark)
    s=blank(prs); bg(s,prs,INK); acc(s,prs,GOLD); eyebrow(s,"The 2026 reset",GOLD)
    text(s,"DMO 8: the rules changed",0.7,0.85,12,0.8,32,"FFFFFF",True)
    rc=[GOLD,TEAL,"5DCAA5",INKMUTE]
    for i,r in enumerate(d["reforms"]):
        cx=0.7+(i%2)*6.15; cyy=2.3+(i//2)*2.0; rect(s,cx,cyy,5.8,1.7,INK2); rect(s,cx,cyy,0.09,1.7,rc[i%4])
        text(s,r[0],cx+0.3,cyy+0.2,5.3,0.5,16,"FFFFFF",True); text(s,r[1],cx+0.3,cyy+0.75,5.3,0.8,12.5,"D9DEE7")
    srcfoot(s,prs,"AER DMO 8 final determination, 26 May 2026",dark=True)

    # 14 YOUR ANALYSIS
    add_custom_analysis(prs,d)

    # 15 CLOSE (dark)
    s=blank(prs); bg(s,prs,INK); acc(s,prs,TEAL,0.18); eyebrow(s,b["org"],TEAL,0.95,0.85)
    text(s,b["tagline"],0.95,1.5,11.5,1.4,38,"FFFFFF",False,DISPLAY)
    text(s,"This entire analysis was generated from public AER determinations - and the tool that builds it is open source.",0.97,3.1,11,0.7,16,"D9DEE7")
    rect(s,0.95,3.95,11.4,1.7,INK2); rect(s,0.95,3.95,0.09,1.7,GOLD)
    text(s,"Clone the open-source AER Market Agent skill, add it to your AI platform (tested on Anthropic's Claude), point it at the latest AER data, customise the skill to your requirements and regenerate this deck in minutes. Ask natural language questions and the skill will analyse.",1.3,4.15,10.7,1.35,15,"FFFFFF",anchor=MSO_ANCHOR.MIDDLE,spacing=1.15)
    text(s,b["url"],0.97,5.95,6,0.4,14,TEAL,True)

    # 16 METHOD
    s=blank(prs); bg(s,prs,PAPER); acc(s,prs,TEAL); eyebrow(s,"Method & sources",TEALD)
    text(s,"How we built this",0.7,0.85,12,0.8,32,NAVY,True)
    meth=[["Data","AER final determinations, DMO 1 (2019-20) to DMO 8 (2026-27)."],
      ["Headline series","Residential flat rate, no controlled load - comparable across all years."],
      ["Cost stack","Wholesale & environmental from AER cost data; network/retail split per AER proportions."],
      ["Network","AER network revenue determinations; AER State of the Energy Market (revenue, asset base, capex, utilisation, tariffs); AEMO 2024 system plan; AER demand & energy charts."],
      ["Inflation","Real terms use ABS CPI (8 capitals, June quarter); 2026-27 estimated."],
      ["Caveat","Small-business benchmark changed 20,000 to 10,000 kWh at DMO 4 - not comparable across it."]]
    for i,m in enumerate(meth):
        my=1.95+i*0.7; text(s,m[0],0.7,my,2.3,0.5,12,TEALD,True); text(s,m[1],3.1,my,5.0,0.6,11,NAVY,spacing=1.05)
    text(s,"Open the sources",8.8,1.9,4,0.3,14,NAVY,True)
    for i,r in enumerate(d["references"]):
        text(s,r[0],8.8,2.35+i*0.56,4.3,0.45,11,TEALD,anchor=MSO_ANCHOR.MIDDLE,link=r[1])

    prs.save(out); return out

# ============================================================
#  ADD YOUR OWN ANALYSIS HERE - this slide is yours.
#  Use text()/rect()/oval() and the chart_*() helpers above.
# ============================================================
def add_custom_analysis(prs,d):
    b=d["branding"]; s=blank(prs); bg(s,prs,b["paper"]); acc(s,prs,b["teal"]); eyebrow(s,"Your analysis",b["teald"])
    text(s,"Add your own slide here",0.7,0.85,12,0.8,32,b["text"],True)
    rect(s,0.7,2.2,12.0,4.0,"FFFFFF",line=b["line"])
    text(s,"This placeholder is the fork point. Edit add_custom_analysis() in scripts/build_report.py to drop in your own chart, forecast or client view - then delete this note.",1.1,2.7,11.2,2.0,16,b["muted"],anchor=MSO_ANCHOR.TOP,spacing=1.3)

def main():
    ap=argparse.ArgumentParser(); ap.add_argument("--dataset",default=os.path.join(HERE,"..","data","dmo_dataset.json")); ap.add_argument("--out",default="DMO Report.pptx")
    a=ap.parse_args(); d=json.load(open(a.dataset,encoding="utf-8")); print("Built:",build(d,a.out))
if __name__=="__main__": main()
